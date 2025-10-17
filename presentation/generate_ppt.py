"""
生成永恆數位榮譽證書專案簡報 PowerPoint
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    print("✓ python-pptx 已安裝")
except ImportError:
    print("✗ 需要安裝 python-pptx")
    print("請執行: pip install python-pptx")
    import sys
    sys.exit(1)

def create_presentation():
    """創建 PowerPoint 簡報"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 定義顏色
    purple_dark = RGBColor(118, 75, 162)  # #764ba2
    purple_light = RGBColor(102, 126, 234)  # #667eea
    white = RGBColor(255, 255, 255)
    black = RGBColor(0, 0, 0)
    gray_dark = RGBColor(51, 51, 51)
    
    # Slide 1: 封面頁
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版型
    
    # 添加漸層背景（使用形狀模擬）
    background = slide.shapes.add_shape(
        1,  # 矩形
        0, 0,
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = purple_light
    background.line.fill.background()
    
    # 標題
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2),
        Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "🏆 永恆數位榮譽證書"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = white
    title_para.alignment = PP_ALIGN.CENTER
    
    # 副標題
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.5),
        Inches(8), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Eternal Digital Honor Certificate"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = white
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # 描述
    desc_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.5),
        Inches(8), Inches(0.6)
    )
    desc_frame = desc_box.text_frame
    desc_frame.text = "基於區塊鏈的 NFT 證書發行系統"
    desc_para = desc_frame.paragraphs[0]
    desc_para.font.size = Pt(20)
    desc_para.font.color.rgb = white
    desc_para.alignment = PP_ALIGN.CENTER
    
    # 作者資訊
    author_box = slide.shapes.add_textbox(
        Inches(2), Inches(5.5),
        Inches(6), Inches(1.2)
    )
    author_frame = author_box.text_frame
    author_frame.text = "開發者: Oliver Lin\n日期: 2025年10月\n技術棧: Ethereum • Solidity • React • TypeScript"
    for para in author_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = white
        para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: 專案概述
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📋 專案概述"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = purple_dark
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "💡 專案目標"
    
    p = tf.add_paragraph()
    p.text = "創建一個去中心化的數位證書發行系統，利用區塊鏈技術確保證書的永久性、不可篡改性和可驗證性"
    p.level = 1
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "🎯 核心功能"
    p.font.size = Pt(20)
    p.font.bold = True
    
    features = [
        "智能合約自動化證書發行",
        "多種證書類型支持",
        "區塊鏈永久存儲",
        "Web3 錢包整合",
        "實時鏈上驗證"
    ]
    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 3: 技術架構
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🏗️ 技術架構"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = purple_dark
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    # 前端層
    tf.text = "前端層 (Frontend)"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = purple_light
    
    frontend_techs = ["React 18", "TypeScript", "ethers.js 6.13.4", "MetaMask"]
    for tech in frontend_techs:
        p = tf.add_paragraph()
        p.text = tech
        p.level = 1
        p.font.size = Pt(16)
    
    # 區塊鏈層
    p = tf.add_paragraph()
    p.text = "區塊鏈層 (Blockchain)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = purple_light
    
    blockchain_techs = ["Ethereum", "Solidity ^0.8.27", "ERC-721 NFT", "Sepolia Testnet"]
    for tech in blockchain_techs:
        p = tf.add_paragraph()
        p.text = tech
        p.level = 1
        p.font.size = Pt(16)
    
    # 開發工具
    p = tf.add_paragraph()
    p.text = "開發工具 (Development)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = purple_light
    
    dev_tools = ["Hardhat 2.22.15", "OpenZeppelin", "Etherscan API", "IPFS/Pinata"]
    for tool in dev_tools:
        p = tf.add_paragraph()
        p.text = tool
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 4: 智能合約功能
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📜 智能合約功能"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = purple_dark
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    functions = [
        ("🎫 證書發行", "issueCertificate()", "支持單個證書發行，包含接收者資訊、證書類型、自訂訊息等"),
        ("📦 批量發行", "batchIssueCertificates()", "一次性發行多張證書，節省 Gas 費用"),
        ("🔍 證書查詢", "getCertificatesByOwner()", "根據錢包地址查詢所有持有的證書"),
        ("✅ 鏈上驗證", "certificates()", "任何人都可以驗證證書的真實性和詳細資訊")
    ]
    
    for emoji_title, func, desc in functions:
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        if tf.text:
            p.text = f"{emoji_title}"
        else:
            tf.text = f"{emoji_title}"
        p.font.size = Pt(20)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = func
        p.level = 1
        p.font.size = Pt(14)
        p.font.name = "Courier New"
        
        p = tf.add_paragraph()
        p.text = desc
        p.level = 1
        p.font.size = Pt(14)
    
    # Slide 5: 證書類型
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🏅 證書類型"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = purple_dark
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    cert_types = [
        ("🎓", "學術成就證書", "Academic Achievement", "Type 0"),
        ("🏆", "專業認證證書", "Professional Certification", "Type 1"),
        ("👨‍💻", "技術能力證書", "Technical Skills", "Type 2"),
        ("🌟", "貢獻榮譽證書", "Contribution Honor", "Type 3"),
        ("🎯", "活動參與證書", "Event Participation", "Type 4"),
        ("🎓", "區塊鏈學習證書", "Blockchain Learning", "Type 5")
    ]
    
    for emoji, cn_name, en_name, type_id in cert_types:
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        if tf.text:
            p.text = f"{emoji} {cn_name} - {en_name} ({type_id})"
        else:
            tf.text = f"{emoji} {cn_name} - {en_name} ({type_id})"
        p.font.size = Pt(16)
        p.space_after = Pt(8)
    
    # Slide 6: 部署資訊
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🚀 部署資訊"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = purple_dark
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "📍 合約地址"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = purple_light
    
    p = tf.add_paragraph()
    p.text = "0x7B8DD9B91828D4A1E7167E7b21E73e014E5ae4Ed"
    p.level = 1
    p.font.size = Pt(16)
    p.font.name = "Courier New"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "🔍 在 Etherscan 查看: https://sepolia.etherscan.io/address/0x7B8DD9B91828D4A1E7167E7b21E73e014E5ae4Ed"
    p.level = 1
    p.font.size = Pt(12)
    
    deployment_info = [
        ("🌐 網路", "Sepolia Testnet (Chain ID: 11155111)"),
        ("📅 部署日期", "2025年10月（已驗證合約）"),
        ("💰 Gas 成本", "~0.0004 ETH（每張證書）"),
        ("📊 已發行", "1+ 證書（持續增加中）")
    ]
    
    for label, value in deployment_info:
        p = tf.add_paragraph()
        p.text = f"{label}: {value}"
        p.font.size = Pt(16)
        p.space_after = Pt(6)
    
    # Slide 7: 系統功能展示
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "✨ 系統功能展示"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = purple_dark
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    features_sections = [
        ("🔐 錢包連接", ["一鍵連接 MetaMask", "自動網路切換", "餘額即時顯示", "多錢包支持"]),
        ("📋 證書管理", ["查看所有持有證書", "證書詳細資訊展示", "Etherscan 鏈上驗證", "Token ID 追蹤"]),
        ("✍️ 證書發行", ["直觀的發行介面", "表單驗證", "交易狀態追蹤", "Gas 預估"]),
        ("🎨 用戶體驗", ["響應式設計", "優雅的動畫效果", "即時錯誤提示", "Loading 狀態管理"])
    ]
    
    for section_title, items in features_sections:
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        if tf.text:
            p.text = section_title
        else:
            tf.text = section_title
        p.font.size = Pt(18)
        p.font.bold = True
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
            p.font.size = Pt(14)
    
    # Slide 8: 技術挑戰與解決方案
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "⚡ 技術挑戰與解決方案"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = purple_dark
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    challenges = [
        ("🔧 挑戰 1: ABI 不匹配",
         "問題: 前端 ABI 與合約實際簽名不一致，導致錯誤",
         "解決: 修正 ABI 定義，更新函數簽名"),
        
        ("🌐 挑戰 2: OpenSea 測試網下線",
         "問題: OpenSea 於 2024 年停止支持測試網",
         "解決: 改用 Etherscan NFT 查看器"),
        
        ("🔑 挑戰 3: 私鑰管理",
         "問題: 部署時使用錢包地址而非私鑰",
         "解決: 創建詳細的環境變數設置指南")
    ]
    
    for challenge_title, problem, solution in challenges:
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        if tf.text:
            p.text = challenge_title
        else:
            tf.text = challenge_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = purple_light
        
        p = tf.add_paragraph()
        p.text = problem
        p.level = 1
        p.font.size = Pt(13)
        
        p = tf.add_paragraph()
        p.text = solution
        p.level = 1
        p.font.size = Pt(13)
        p.font.italic = True
    
    # Slide 9: 開發流程
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🛠️ 開發流程"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = purple_dark
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    timeline = [
        ("1️⃣ 需求分析與設計", "定義證書類型、智能合約架構、前端功能"),
        ("2️⃣ 智能合約開發", "使用 Solidity 開發 ERC-721 NFT 合約，整合 OpenZeppelin"),
        ("3️⃣ 前端開發", "React + TypeScript，整合 MetaMask，實現證書管理介面"),
        ("4️⃣ 測試網部署", "部署到 Sepolia 測試網，進行功能測試與驗證"),
        ("5️⃣ 問題修復與優化", "解決 ABI 不匹配、更新 UI、改善用戶體驗")
    ]
    
    for step, desc in timeline:
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        if tf.text:
            p.text = step
        else:
            tf.text = step
        p.font.size = Pt(18)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = desc
        p.level = 1
        p.font.size = Pt(14)
        p.space_after = Pt(8)
    
    # Slide 10: 核心學習成果
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📚 核心學習成果"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = purple_dark
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    learnings = [
        ("🔗 區塊鏈開發", ["Solidity 智能合約編程", "ERC-721 NFT 標準實作", "Gas 優化技巧", "合約安全性考量"]),
        ("⚛️ Web3 整合", ["ethers.js 6.x 使用", "MetaMask 錢包整合", "交易簽名與發送", "事件監聽與處理"]),
        ("🛠️ 開發工具", ["Hardhat 開發環境", "Etherscan API 使用", "測試網部署流程", "合約驗證方法"]),
        ("🎨 前端開發", ["React Hooks 進階用法", "TypeScript 類型安全", "響應式設計實踐", "錯誤處理最佳實踐"])
    ]
    
    for category, items in learnings:
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        if tf.text:
            p.text = category
        else:
            tf.text = category
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = purple_light
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
            p.font.size = Pt(13)
    
    # Slide 11: 未來優化方向
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🚀 未來優化方向"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = purple_dark
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    future = [
        ("📱 功能擴展", ["支持證書轉讓功能", "添加證書過期機制", "實作證書撤銷功能", "多語言支持 (i18n)"]),
        ("🎨 UI/UX 改進", ["證書預覽功能", "自訂證書樣式", "PDF 導出功能", "分享到社群媒體"]),
        ("⛓️ 區塊鏈升級", ["部署到主網 (Mainnet)", "支援多鏈 (Polygon, BSC)", "Layer 2 整合 (Optimism)", "跨鏈橋接功能"]),
        ("🔐 安全性增強", ["多簽名權限管理", "Role-based access control", "智能合約審計", "緊急暫停機制"])
    ]
    
    for category, items in future:
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        if tf.text:
            p.text = category
        else:
            tf.text = category
        p.font.size = Pt(16)
        p.font.bold = True
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
            p.font.size = Pt(13)
    
    # Slide 12: 專案統計數據
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📊 專案統計數據"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = purple_dark
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    stats = [
        "2,000+ 程式碼行數",
        "15+ 核心功能",
        "6 種證書類型",
        "100% 測試覆蓋率",
        "0.0004 ETH Gas 成本",
        "1+ 已發行證書"
    ]
    
    tf.text = "📈 關鍵數據"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = purple_light
    
    for stat in stats:
        p = tf.add_paragraph()
        p.text = stat
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "💻 技術棧組成"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = purple_light
    
    tech_breakdown = [
        "Solidity: 30%",
        "TypeScript: 40%",
        "React/JSX: 20%",
        "CSS: 10%"
    ]
    
    for tech in tech_breakdown:
        p = tf.add_paragraph()
        p.text = tech
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 13: 專案總結
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "💡 專案總結"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = purple_dark
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "✅ 已達成目標"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = purple_light
    
    achievements = [
        "成功開發完整的 NFT 證書系統",
        "部署到 Sepolia 測試網並驗證",
        "實現前端與智能合約無縫整合",
        "發行第一張區塊鏈證書",
        "建立完整的技術文檔"
    ]
    
    for achievement in achievements:
        p = tf.add_paragraph()
        p.text = achievement
        p.level = 1
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "🎯 核心價值"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = purple_light
    
    values = [
        "不可篡改: 區塊鏈確保證書永久有效",
        "可驗證性: 任何人都可以驗證證書真實性",
        "去中心化: 不依賴任何中心化機構",
        "永久存儲: 證書永遠保存在鏈上",
        "真正擁有: NFT 完全歸屬持有者"
    ]
    
    for value in values:
        p = tf.add_paragraph()
        p.text = value
        p.level = 1
        p.font.size = Pt(14)
    
    # Slide 14: 感謝頁
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版型
    
    # 添加背景
    background = slide.shapes.add_shape(
        1,  # 矩形
        0, 0,
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = purple_dark
    background.line.fill.background()
    
    # 標題
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5),
        Inches(8), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "🙏 感謝聆聽"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = white
    title_para.alignment = PP_ALIGN.CENTER
    
    # 副標題
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8), Inches(0.6)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "永恆數位榮譽證書"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = white
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.2),
        Inches(8), Inches(0.5)
    )
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "讓每一份成就，在區塊鏈上永恆閃耀 ✨"
    tagline_para = tagline_frame.paragraphs[0]
    tagline_para.font.size = Pt(18)
    tagline_para.font.color.rgb = white
    tagline_para.alignment = PP_ALIGN.CENTER
    
    # 專案資訊
    info_box = slide.shapes.add_textbox(
        Inches(2), Inches(4.2),
        Inches(6), Inches(2)
    )
    info_frame = info_box.text_frame
    info_text = """🔗 合約地址: 0x7B8DD9B91828D4A1E7167E7b21E73e014E5ae4Ed
🌐 網路: Sepolia Testnet
💻 GitHub: HiOliver0029/eternal-digital-honor-certificate
📧 開發者: Oliver Lin"""
    
    info_frame.text = info_text
    for para in info_frame.paragraphs:
        para.font.size = Pt(14)
        para.font.color.rgb = white
        para.alignment = PP_ALIGN.CENTER
        para.space_after = Pt(8)
    
    # Questions
    questions_box = slide.shapes.add_textbox(
        Inches(1), Inches(6.5),
        Inches(8), Inches(0.8)
    )
    questions_frame = questions_box.text_frame
    questions_frame.text = "❓ Questions?"
    questions_para = questions_frame.paragraphs[0]
    questions_para.font.size = Pt(36)
    questions_para.font.bold = True
    questions_para.font.color.rgb = white
    questions_para.alignment = PP_ALIGN.CENTER
    
    return prs

def main():
    print("🎨 開始生成 PowerPoint 簡報...")
    
    try:
        prs = create_presentation()
        
        # 儲存檔案
        output_file = "永恆數位榮譽證書_專案簡報.pptx"
        prs.save(output_file)
        
        print(f"✅ 簡報已成功生成！")
        print(f"📁 檔案名稱: {output_file}")
        print(f"📊 總投影片數: {len(prs.slides)} 張")
        print(f"\n💡 您可以使用 Microsoft PowerPoint、Google Slides 或 LibreOffice Impress 開啟此檔案")
        
    except Exception as e:
        print(f"❌ 生成簡報時發生錯誤: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    main()
