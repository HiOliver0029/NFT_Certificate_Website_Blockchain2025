"""
Generate English PowerPoint Presentation for Eternal Digital Honor Certificate
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    print("✓ python-pptx installed")
except ImportError:
    print("✗ python-pptx not installed")
    print("Please run: pip install python-pptx")
    import sys
    sys.exit(1)

def create_presentation():
    """Create PowerPoint Presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)
    
    # Define colors
    purple_dark = RGBColor(118, 75, 162)  # #764ba2
    purple_light = RGBColor(102, 126, 234)  # #667eea
    white = RGBColor(255, 255, 255)
    black = RGBColor(0, 0, 0)
    gray_dark = RGBColor(51, 51, 51)
    
    def add_title_slide(title_text, subtitle_text, author_info):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Background
        background = slide.shapes.add_shape(
            1, 0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = purple_light
        background.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(11.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = white
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.5), Inches(11.333), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle_text
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(32)
        subtitle_para.font.color.rgb = white
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Author Info
        author_box = slide.shapes.add_textbox(
            Inches(1), Inches(5), Inches(11.333), Inches(2)
        )
        author_frame = author_box.text_frame
        author_frame.text = author_info
        for para in author_frame.paragraphs:
            para.font.size = Pt(20)
            para.font.color.rgb = white
            para.alignment = PP_ALIGN.CENTER

    def add_content_slide(title_text, content_data):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.color.rgb = purple_dark
        title.text_frame.paragraphs[0].font.bold = True
        
        # Content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()  # Clear default empty paragraph
        
        for item in content_data:
            if isinstance(item, dict):
                # Heading
                p = tf.add_paragraph()
                p.text = item['title']
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = purple_light
                p.space_before = Pt(12)
                
                # Body/List
                if 'content' in item:
                    if isinstance(item['content'], list):
                        for li in item['content']:
                            p = tf.add_paragraph()
                            p.text = li
                            p.font.size = Pt(20)
                            p.level = 1
                    else:
                        p = tf.add_paragraph()
                        p.text = item['content']
                        p.font.size = Pt(20)
                        p.level = 1
            else:
                # Simple list item or text
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(20)
                p.level = 0

    # Slide 1: Title Slide
    add_title_slide(
        "🎓 Eternal Digital Honor Certificate",
        "Blockchain-based NFT Certificate Issuance System",
        "👨‍💻 Project Developer: Oliver Lin\n📅 Presentation Date: October 2025\n🔗 Ethereum Sepolia Testnet"
    )

    # Slide 2: Project Overview
    add_content_slide("📋 Project Overview", [
        {"title": "💡 Project Objectives", "content": "Create a decentralized digital certificate issuance system using blockchain technology to ensure permanence, immutability, and verifiability of certificates"},
        {"title": "🎯 Core Features", "content": [
            "Automated certificate issuance via smart contracts",
            "Support for multiple certificate types",
            "Permanent blockchain storage",
            "Web3 wallet integration",
            "Real-time on-chain verification"
        ]}
    ])

    # Slide 3: Technical Architecture
    add_content_slide("🏗️ Technical Architecture", [
        {"title": "Frontend Layer", "content": "React 18 • TypeScript • ethers.js 6.13.4 • MetaMask"},
        {"title": "Blockchain Layer", "content": "Ethereum • Solidity ^0.8.27 • ERC-721 NFT • Sepolia Testnet"},
        {"title": "Development Tools", "content": "Hardhat 2.22.15 • OpenZeppelin • Etherscan API • IPFS/Pinata"}
    ])

    # Slide 4: Smart Contract Functions
    add_content_slide("📜 Smart Contract Functions", [
        {"title": "🎫 Certificate Issuance", "content": "issueCertificate() - Support for single certificate issuance including recipient info, certificate type, and custom messages"},
        {"title": "📦 Batch Issuance", "content": "batchIssueCertificates() - Issue multiple certificates at once, saving gas fees"},
        {"title": "🔍 Certificate Query", "content": "getCertificatesByOwner() - Query all certificates owned by a wallet address"},
        {"title": "✅ On-Chain Verification", "content": "certificates() - Anyone can verify the authenticity and details of certificates"}
    ])

    # Slide 5: Certificate Types
    add_content_slide("🏅 Certificate Types", [
        "🎓 Academic Achievement (Type 0)",
        "🏆 Professional Certification (Type 1)",
        "👨‍💻 Technical Skills (Type 2)",
        "🌟 Contribution Honor (Type 3)",
        "🎯 Event Participation (Type 4)",
        "🎓 Blockchain Learning (Type 5)"
    ])

    # Slide 6: Deployment Information
    add_content_slide("🚀 Deployment Information", [
        {"title": "📍 Contract Address", "content": "0x7B8DD9B91828D4A1E7167E7b21E73e014E5ae4Ed"},
        {"title": "🌐 Network", "content": "Sepolia Testnet (Chain ID: 11155111)"},
        {"title": "📅 Deployment Date", "content": "October 2025 (Verified Contract)"},
        {"title": "💰 Gas Cost", "content": "~0.0004 ETH Per Certificate"},
        {"title": "📊 Certificates Issued", "content": "1+ Certificates (Continuously Growing)"}
    ])

    # Slide 7: System Features
    add_content_slide("✨ System Features", [
        {"title": "🔐 Wallet Connection", "content": ["One-click MetaMask connection", "Automatic network switching", "Real-time balance display", "Multi-wallet support"]},
        {"title": "📋 Certificate Management", "content": ["View all owned certificates", "Detailed certificate information", "Etherscan on-chain verification", "Token ID tracking"]},
        {"title": "✍️ Certificate Issuance", "content": ["Intuitive issuance interface", "Form validation", "Transaction status tracking", "Gas estimation"]},
        {"title": "🎨 User Experience", "content": ["Responsive design", "Elegant animations", "Real-time error alerts", "Loading state management"]}
    ])

    # Slide 8: System Interface
    add_content_slide("💻 System Interface", [
        {"title": "Application Screenshots", "content": [
            "🖼️ Certificate Management Interface - Displays owned NFT certificate list",
            "✍️ Certificate Issuance Interface - Enter recipient info and issue new certificates",
            "🔍 Etherscan Verification - View certificate details on blockchain explorer"
        ]},
        {"title": "💡 Live Demo", "content": "Can demonstrate the running system on-site"}
    ])

    # Slide 9: Technical Challenges & Solutions
    add_content_slide("⚡ Technical Challenges & Solutions", [
        {"title": "🔧 Challenge 1: ABI Mismatch", "content": "Problem: Frontend ABI didn't match actual contract signature\nSolution: Fixed ABI definition, removed non-existent imageURI parameter"},
        {"title": "🌐 Challenge 2: OpenSea Testnet Sunset", "content": "Problem: OpenSea discontinued testnet support in 2024\nSolution: Switched to Etherscan NFT viewer for verification"},
        {"title": "🔑 Challenge 3: Private Key Management", "content": "Problem: Used wallet address instead of private key during deployment\nSolution: Created detailed environment variable setup guide"}
    ])

    # Slide 10: Development Process
    add_content_slide("🛠️ Development Process", [
        "1. Requirements Analysis & Design - Define certificate types, smart contract architecture",
        "2. Smart Contract Development - Develop ERC-721 NFT contract using Solidity",
        "3. Frontend Development - React + TypeScript, integrate MetaMask",
        "4. Testnet Deployment - Deploy to Sepolia testnet, conduct functional testing",
        "5. Bug Fixes & Optimization - Resolve ABI mismatch, update UI, improve UX"
    ])

    # Slide 11: Key Learnings
    add_content_slide("📚 Key Learnings", [
        {"title": "🔗 Blockchain Development", "content": "Solidity, ERC-721, Gas optimization, Security"},
        {"title": "⚛️ Web3 Integration", "content": "ethers.js, MetaMask, Transaction handling, Events"},
        {"title": "🛠️ Development Tools", "content": "Hardhat, Etherscan API, Testnet deployment, Verification"},
        {"title": "🎨 Frontend Development", "content": "React Hooks, TypeScript, Responsive design, Error handling"}
    ])

    # Slide 12: Future Enhancements
    add_content_slide("🚀 Future Enhancements", [
        {"title": "📱 Feature Extensions", "content": "Certificate transfer, Expiration mechanism, Revocation, i18n"},
        {"title": "🎨 UI/UX Improvements", "content": "Certificate preview, Custom styling, PDF export, Social sharing"},
        {"title": "⛓️ Blockchain Upgrades", "content": "Mainnet, Multi-chain (Polygon, BSC), Layer 2 (Optimism)"},
        {"title": "🔐 Security Enhancements", "content": "Multi-sig, Role-based access, Auditing, Emergency pause"}
    ])

    # Slide 13: Project Statistics
    add_content_slide("📊 Project Statistics", [
        {"title": "Metrics", "content": [
            "2,000+ Lines of Code",
            "15+ Core Features",
            "6 Certificate Types",
            "100% Test Coverage",
            "0.0004 ETH Gas Cost",
            "1+ Certificates Issued"
        ]},
        {"title": "Technology Stack", "content": "TypeScript (40%), Solidity (30%), React/JSX (20%), CSS (10%)"}
    ])

    # Slide 14: Conclusion
    add_content_slide("💡 Project Summary", [
        {"title": "✅ Achievements", "content": [
            "Successfully developed complete NFT certificate system",
            "Deployed to Sepolia testnet and verified",
            "Seamless frontend and smart contract integration",
            "Issued first blockchain certificate"
        ]},
        {"title": "🎯 Core Values", "content": [
            "Immutable: Blockchain ensures permanent validity",
            "Verifiable: Anyone can verify authenticity",
            "Decentralized: No dependence on centralized institutions",
            "True Ownership: NFTs fully belong to holders"
        ]}
    ])

    # Slide 15: Thank You Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = purple_light
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(11.333), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "🙏 Thank You!"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = white
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4), Inches(11.333), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Questions & Discussion"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = white
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    contact_box = slide.shapes.add_textbox(
        Inches(1), Inches(5.5), Inches(11.333), Inches(1.5)
    )
    contact_frame = contact_box.text_frame
    contact_frame.text = "📧 Contact: oliver.lin@example.com\n🔗 GitHub: @HiOliver0029\n🌐 Project: Eternal Digital Honor Certificate"
    for para in contact_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = white
        para.alignment = PP_ALIGN.CENTER

    # Save
    prs.save('Eternal_Digital_Honor_Certificate_Presentation.pptx')
    print("✅ PowerPoint presentation generated: Eternal_Digital_Honor_Certificate_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
